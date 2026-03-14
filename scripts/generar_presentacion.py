# -*- coding: utf-8 -*-
"""
Genera la presentación PowerPoint del proyecto Gestión Humana – Colbeef.
Ejecutar: python scripts/generar_presentacion.py
Requiere: pip install python-pptx
"""
import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    try:
        from pptx.dml.color import RgbColor
        VERDE_OSCURO = RgbColor(0x0B, 0x35, 0x18)
        VERDE = RgbColor(0x2D, 0x9E, 0x3F)
        GRIS_TEXTO = RgbColor(0x37, 0x41, 0x51)
    except ImportError:
        VERDE_OSCURO = VERDE = GRIS_TEXTO = None
except ImportError as e:
    print("Instale python-pptx: pip install python-pptx")
    print("Error:", e)
    sys.exit(1)


def set_title_style(tf):
    """Ajusta fuente del título."""
    for p in tf.paragraphs:
        p.font.size = Pt(28)
        p.font.bold = True
        if VERDE_OSCURO and hasattr(p.font, 'color') and p.font.color:
            try:
                p.font.color.rgb = VERDE_OSCURO
            except Exception:
                pass


def set_body_style(tf, font_size=14):
    """Ajusta fuente del cuerpo."""
    for p in tf.paragraphs:
        p.font.size = Pt(font_size)
        if GRIS_TEXTO:
            try:
                p.font.color.rgb = GRIS_TEXTO
            except Exception:
                pass


def add_slide_title_content(prs, title_text, bullets):
    """Diapositiva título + viñetas."""
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title_text
    set_title_style(slide.shapes.title.text_frame)
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.level = 0
        p.font.size = Pt(14)
        try:
            p.font.color.rgb = GRIS_TEXTO
        except Exception:
            pass
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ─── Portada ───
    layout0 = prs.slide_layouts[0]  # Title
    slide = prs.slides.add_slide(layout0)
    slide.shapes.title.text = "Sistema de Gestión Humana"
    slide.placeholders[1].text = "Colbeef · Python Flask & MySQL"
    set_title_style(slide.shapes.title.text_frame)
    if VERDE:
        try:
            for p in slide.placeholders[1].text_frame.paragraphs:
                p.font.size = Pt(18)
                p.font.color.rgb = VERDE
        except Exception:
            pass

    # ─── Introducción ───
    add_slide_title_content(prs, "Introducción", [
        "Sistema web interno para la gestión de recursos humanos de Colbeef.",
        "Centraliza: personal, organización, familia, eventos, permisos y vacaciones.",
        "Múltiples roles (Coordinación GH, Gestor Contratación, Bienestar, Nómina, SST, Admin).",
        "Portal del empleado para solicitar permisos y consultar sus solicitudes.",
        "Notificaciones por correo y firma digital sobre el formato de permiso aprobado.",
    ])

    # ─── Objetivos ───
    add_slide_title_content(prs, "Objetivos", [
        "Digitalizar procesos de Gestión Humana (permisos, licencias, vacaciones).",
        "Dar visibilidad a la coordinación sobre solicitudes y evidencia adjunta.",
        "Notificar por correo a coordinación, gestor y empleado en cada paso.",
        "Mantener trazabilidad (auditoría, quién aprobó/rechazó y cuándo).",
        "Permitir consulta por parte del empleado (mis solicitudes, estado).",
    ])

    # ─── Tecnologías ───
    add_slide_title_content(prs, "Tecnologías", [
        "Backend: Python 3.10+, Flask.",
        "Base de datos: MySQL (gestio_humana).",
        "Frontend: HTML/CSS/JS, plantillas Jinja2.",
        "Correo: SMTP (mail Colbeef) para notificaciones.",
        "PDF: pypdf, reportlab, Pillow (firma digital sobre formato de permiso).",
        "Despliegue: Waitress (producción) o servidor de desarrollo Flask.",
    ])

    # ─── Módulos principales ───
    add_slide_title_content(prs, "Módulos principales", [
        "Organización: departamentos, áreas, perfiles ocupacionales.",
        "Personal: empleados activos/inactivos, fichas, retiro.",
        "Familia: hijos activos/inactivos.",
        "Eventos: cumpleaños, aniversario laboral.",
        "EPS y fondos de pensiones.",
        "Reportes y exportación a Excel.",
        "Permisos y licencias: solicitud, aprobación, evidencia, firma.",
        "Vacaciones: solicitud y resolución.",
        "Admin: usuarios, roles, módulos visibles por rol.",
    ])

    # ─── Flujo de permisos ───
    add_slide_title_content(prs, "Flujo de permisos / licencias", [
        "Empleado (o gestor en su nombre) llena el formato GH-FR-007 en la web.",
        "Si es no remunerado: adjunta evidencia (PDF o imagen).",
        "Correo a Coordinación GH y a Gestión Humana (informativa) y a Gestor de Contratación.",
        "Coordinación revisa la solicitud y la evidencia (enlace «Ver evidencia» en la tarjeta).",
        "Coordinación aprueba o rechaza; opcionalmente agrega observaciones.",
        "Empleado recibe correo con el resultado; si hay evidencia PDF y firma configurada, se adjunta copia firmada.",
    ])

    # ─── Roles ───
    add_slide_title_content(prs, "Roles y permisos", [
        "ADMIN: acceso total, catálogos, usuarios.",
        "COORD. GH: mismo nivel que admin en operación; aprueba/rechaza permisos y vacaciones.",
        "GH INFORMADA: ve permisos y recibe notificaciones; no aprueba ni rechaza.",
        "GESTOR DE CONTRATACIÓN: personal, permisos (consulta y solicitudes).",
        "BIENESTAR SOCIAL, GESTOR DE NÓMINA, GESTOR SST: según módulos asignados (lectura/escritura).",
        "EMPLEADO: portal con solicitud de permiso y mis solicitudes.",
    ])

    # ─── Seguridad y buenas prácticas ───
    add_slide_title_content(prs, "Seguridad y buenas prácticas", [
        "Autenticación por sesión; cambio de clave obligatorio si se fuerza.",
        "Variables sensibles en .env (no en el repositorio): MySQL, SMTP, SECRET_KEY.",
        "Evidencia y archivos en instance/uploads; descarga solo con permiso del módulo.",
        "Registro de auditoría para acciones relevantes (aprobaciones, rechazos).",
    ])

    # ─── Cierre ───
    add_slide_title_content(prs, "Resumen", [
        "Sistema listo para uso en Colbeef: permisos, notificaciones, evidencia y firma digital.",
        "Escalable a más módulos y roles según necesidad.",
        "Documentación e instalación en README y scripts de base de datos en /database.",
    ])

    # ─── Gracias ───
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Gracias"
    set_title_style(slide.shapes.title.text_frame)
    body = slide.placeholders[1]
    body.text_frame.text = "Sistema de Gestión Humana – Colbeef"
    set_body_style(body.text_frame, font_size=16)

    # Guardar
    out_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    out_path = os.path.join(out_dir, "presentacion_gestio_humana.pptx")
    prs.save(out_path)
    print(f"Presentación generada: {out_path}")
    return out_path


if __name__ == "__main__":
    main()
